#!/usr/bin/env python3
"""
export-native-pptx-v4.py — 改进版 v4

v4 修复：
1. extractSegments 追踪 fontWeight（bold）→ <strong>/<b> 标签在 PPTX 中正确加粗
2. 修复卡片边框处理：区分全边框 vs 单侧边框
3. 合并逻辑同时考虑 color 和 fontWeight
"""

import sys
import argparse
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
import re

def check_deps():
    missing = []
    try:
        from playwright.sync_api import sync_playwright
    except ImportError:
        missing.append("playwright")
    try:
        from pptx import Presentation
    except ImportError:
        missing.append("python-pptx")
    if missing:
        print(f"Install: pip install {' '.join(missing)}")
        sys.exit(1)

check_deps()

from playwright.sync_api import sync_playwright, Page
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


def set_roundrect_adj(shape, radius_px: float, width_in: float, height_in: float):
    """Set rounded rectangle corner radius via OOXML adj value."""
    from lxml import etree
    NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    radius_in = radius_px / 96.0
    shorter = min(width_in, height_in)
    if shorter <= 0:
        return
    adj = int(radius_in / (shorter / 2) * 100000)
    adj = max(0, min(50000, adj))
    prstGeom = shape._element.spPr.find(f'{{{NS}}}prstGeom')
    if prstGeom is None:
        return
    avLst = prstGeom.find(f'{{{NS}}}avLst')
    if avLst is None:
        avLst = etree.SubElement(prstGeom, f'{{{NS}}}avLst')
    for gd in avLst.findall(f'{{{NS}}}gd'):
        avLst.remove(gd)
    gd = etree.SubElement(avLst, f'{{{NS}}}gd')
    gd.set('name', 'adj')
    gd.set('fmla', f'val {adj}')


def suppress_line(shape):
    """Write <a:ln><a:noFill/></a:ln> directly into spPr to remove shape outline."""
    from lxml import etree
    NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    spPr = shape._element.spPr
    ln_tag = f'{{{NS}}}ln'
    ln = spPr.find(ln_tag)
    if ln is not None:
        spPr.remove(ln)
    ln = etree.SubElement(spPr, ln_tag)
    etree.SubElement(ln, f'{{{NS}}}noFill')


def parse_color(css_color: str, bg: Tuple[int,int,int] = (255, 255, 255)) -> Optional[Tuple[int, int, int]]:
    """Parse a CSS color string, blending rgba() alpha over the given bg color (default white)."""
    if not css_color or css_color in ('transparent',) or 'rgba(0, 0, 0, 0)' in css_color:
        return None
    m = re.search(r'rgba?\((\d+),\s*(\d+),\s*(\d+)(?:,\s*([\d.]+))?\)', css_color)
    if m:
        r, g, b = int(m.group(1)), int(m.group(2)), int(m.group(3))
        a = float(m.group(4)) if m.group(4) else 1.0
        if a <= 0:
            return None
        if a < 1.0:
            r = int(a * r + (1 - a) * bg[0])
            g = int(a * g + (1 - a) * bg[1])
            b = int(a * b + (1 - a) * bg[2])
        return (r, g, b)
    m = re.search(r'#([0-9a-fA-F]{6}|[0-9a-fA-F]{3})', css_color)
    if m:
        h = m.group(1)
        if len(h) == 3:
            h = ''.join([c*2 for c in h])
        return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    return None


def px_to_pt(px_value: str) -> float:
    m = re.search(r'([\d.]+)px', str(px_value))
    if m:
        return round(float(m.group(1)) * 0.75, 1)
    return 12.0


def inject_visible(page: Page):
    page.evaluate(
        'document.querySelectorAll(".slide").forEach(s => s.classList.add("visible"));'
        'document.querySelectorAll(".reveal").forEach(el => {'
        '  el.style.opacity = "1"; el.style.transform = "none";'
        '});'
    )


# JS extraction code (stored as a string, careful with escaping)
_EXTRACT_JS = r"""
(slideIndex) => {
    const slide = document.querySelectorAll('.slide')[slideIndex];
    if (!slide) return {background: null, elements: [], slideSize: null};

    // Fallback: if no .slide-content wrapper, use the slide itself as content root
    const content = slide.querySelector('.slide-content') || slide;

    const slideRect = slide.getBoundingClientRect();
    const bodyStyle = window.getComputedStyle(document.body);
    let bgColor = bodyStyle.backgroundColor;
    if (!bgColor || bgColor === 'transparent' || bgColor === 'rgba(0, 0, 0, 0)') {
        // Body may use a CSS gradient (background shorthand) — extract first color stop
        const bgImg = bodyStyle.backgroundImage || '';
        if (bgImg.includes('gradient')) {
            const cm = bgImg.match(/rgba?\([^)]+\)/g);
            bgColor = (cm && cm.length > 0) ? cm[0] : null;
        } else {
            bgColor = null;
        }
    }

    const slideW = slideRect.width;
    const slideH = slideRect.height;

    // Helper: extract inline text segments with colors
    function extractSegments(el) {
        const elStyle = window.getComputedStyle(el);
        const bgImage = elStyle.backgroundImage || '';
        const bgClip = elStyle.webkitBackgroundClip || elStyle.backgroundClip || '';
        const isGradient = bgImage.includes('gradient') && bgClip === 'text';

        let gradientColor = null;
        if (isGradient) {
            const cm = bgImage.match(/rgba?\([^)]+\)/g);
            if (cm) gradientColor = cm[cm.length - 1];
        }

        let gradientColors = null;
        if (isGradient) {
            const cm2 = bgImage.match(/rgba?\([^)]+\)/g);
            if (cm2 && cm2.length >= 2) gradientColors = [cm2[0], cm2[cm2.length-1]];
        }

        const segments = [];
        function walk(node, color, bold, fontSize) {
            if (node.nodeType === 3) {
                const t = node.textContent;
                // Preserve internal whitespace (e.g. " text" after <strong>Note:</strong>)
                // but skip nodes that are purely whitespace
                if (t && t.trim()) segments.push({text: t, color: color, bold: bold, fontSize: fontSize});
            } else if (node.nodeType === 1) {
                const tag = node.tagName;
                if (tag === 'BR') { segments.push({text: '\n', color: color, bold: bold, fontSize: fontSize}); return; }
                const s2 = window.getComputedStyle(node);
                const bi = s2.backgroundImage || '';
                const bc = s2.webkitBackgroundClip || s2.backgroundClip || '';
                let c = color;
                if (bi.includes('gradient') && bc === 'text') {
                    const cm = bi.match(/rgba?\([^)]+\)/g);
                    if (cm) c = cm[cm.length - 1];
                } else {
                    const sc = s2.color;
                    if (sc && sc !== 'rgba(0, 0, 0, 0)') c = sc;
                }
                let b = bold;
                const fw = s2.fontWeight;
                if (fw === 'bold' || fw === '700' || fw === '800' || fw === '900' || parseInt(fw) >= 600) b = true;
                const fs = s2.fontSize || fontSize;
                for (const child of node.childNodes) walk(child, c, b, fs);
            }
        }

        const baseColor = isGradient ? (gradientColor || elStyle.color) : elStyle.color;
        const baseBold = parseInt(elStyle.fontWeight) >= 600;
        const baseFontSize = elStyle.fontSize;
        for (const child of el.childNodes) walk(child, baseColor, baseBold, baseFontSize);

        // Merge consecutive same-color+bold+fontSize segments
        const merged = [];
        for (const seg of segments) {
            if (merged.length > 0 && merged[merged.length-1].color === seg.color &&
                merged[merged.length-1].bold === seg.bold &&
                merged[merged.length-1].fontSize === seg.fontSize &&
                seg.text !== '\n' && merged[merged.length-1].text !== '\n') {
                merged[merged.length-1].text += seg.text;
            } else {
                merged.push({text: seg.text, color: seg.color, bold: seg.bold, fontSize: seg.fontSize});
            }
        }

        return {segments: merged, gradientColors: gradientColors};
    }

    const TEXT_TAGS = new Set(['h1','h2','h3','h4','h5','h6','p','li','span','a']);

    // Flat recursive traversal
    function flatExtract(el) {
        const rect = el.getBoundingClientRect();
        if (rect.width < 1 || rect.height < 1) return [];

        const tag = el.tagName.toLowerCase();
        const style = window.getComputedStyle(el);

        const bounds = {
            x: (rect.left - slideRect.left) / 96,
            y: (rect.top - slideRect.top) / 96,
            width: rect.width / 96,
            height: rect.height / 96
        };

        const bgColor = style.backgroundColor;
        const hasBg = bgColor && bgColor !== 'transparent' && bgColor !== 'rgba(0, 0, 0, 0)';

        // Check all border sides
        const borderLeft = style.borderLeft || '';
        const borderRight = style.borderRight || '';
        const borderTop = style.borderTop || '';
        const borderBottom = style.borderBottom || '';
        const borderStr = style.border || '';

        const hasLeftBorder = borderLeft && !borderLeft.includes('none') && !borderLeft.startsWith('0px');
        const hasRightBorder = borderRight && !borderRight.includes('none') && !borderRight.startsWith('0px');
        const hasTopBorder = borderTop && !borderTop.includes('none') && !borderTop.startsWith('0px');
        const hasBottomBorder = borderBottom && !borderBottom.includes('none') && !borderBottom.startsWith('0px');
        const hasGeneralBorder = borderStr && !borderStr.includes('none') && !borderStr.startsWith('0px');
        const hasBorder = hasLeftBorder || hasRightBorder || hasTopBorder || hasBottomBorder || hasGeneralBorder;

        const results = [];

        if (TEXT_TAGS.has(tag)) {
            // Text element - render as text box
            const {segments, gradientColors} = extractSegments(el);
            const text = el.innerText.trim();
            if (text || segments.length > 0) {
                results.push({
                    type: 'text',
                    tag: tag,
                    text: text,
                    segments: segments,
                    gradientColors: gradientColors,
                    textTransform: style.textTransform,
                    bounds: bounds,
                    styles: {
                        fontSize: style.fontSize,
                        fontWeight: style.fontWeight,
                        fontFamily: style.fontFamily,
                        letterSpacing: style.letterSpacing,
                        color: style.color,
                        textAlign: style.textAlign,
                        lineHeight: style.lineHeight,
                        listStyleType: style.listStyleType
                    }
                });
            }
            // Don't recurse into text elements
            return results;
        }

        if (tag === 'table') {
            // Extract table structure: all rows and cells with bounds/styles
            const tableRows = [];
            const allRows = el.querySelectorAll('tr');
            allRows.forEach((row) => {
                const isHeader = !!row.closest('thead');
                const cells = row.querySelectorAll('th, td');
                const rowCells = [];
                cells.forEach(cell => {
                    const cellRect = cell.getBoundingClientRect();
                    const cellStyle = window.getComputedStyle(cell);
                    const {segments} = extractSegments(cell);
                    rowCells.push({
                        bounds: {
                            x: (cellRect.left - slideRect.left) / 96,
                            y: (cellRect.top - slideRect.top) / 96,
                            width: cellRect.width / 96,
                            height: cellRect.height / 96
                        },
                        text: cell.innerText.trim(),
                        segments: segments,
                        isHeader: isHeader,
                        styles: {
                            fontSize: cellStyle.fontSize,
                            fontWeight: cellStyle.fontWeight,
                            color: cellStyle.color,
                            backgroundColor: cellStyle.backgroundColor,
                            textAlign: cellStyle.textAlign,
                            paddingLeft: cellStyle.paddingLeft,
                            paddingRight: cellStyle.paddingRight,
                            paddingTop: cellStyle.paddingTop,
                            paddingBottom: cellStyle.paddingBottom,
                            fontFamily: cellStyle.fontFamily,
                            letterSpacing: cellStyle.letterSpacing,
                            borderBottom: cellStyle.borderBottom,
                            borderRight: cellStyle.borderRight
                        }
                    });
                });
                if (rowCells.length > 0) tableRows.push({isHeader: isHeader, cells: rowCells});
            });
            results.push({type: 'table', bounds: bounds, rows: tableRows});
            return results;  // Don't recurse into table children
        }

        if (tag === 'div' || tag === 'section' || tag === 'article' || tag === 'ul' || tag === 'ol') {
            // Filter out decorative elements with highly transparent backgrounds and no text
            // (e.g. ambient orb/cloud divs in blue-sky style: rgba(x,x,x,0.3) blobs)
            const alphaMatch = (bgColor || '').match(/rgba\([^)]+,\s*([\d.]+)\s*\)/);
            const bgAlpha = alphaMatch ? parseFloat(alphaMatch[1]) : 1.0;
            if (bgAlpha < 0.5 && !el.innerText.trim() && bounds.width > 1.5 && bounds.height > 1.5) return [];

            // Detect "leaf-text container": a div whose entire visible content is text
            // Case A: no child elements at all — e.g. <div class="chapter-num">01</div>
            // Case B: only inline child elements + sibling text nodes — e.g. <div class="co"><strong>Note:</strong> more text</div>
            const INLINE_TAGS = new Set(['STRONG','EM','B','I','SPAN','A','MARK','CODE','SMALL']);
            const totalText = el.innerText.trim();
            const allChildInline = el.children.length > 0 &&
                Array.from(el.children).every(c => INLINE_TAGS.has(c.tagName));
            const childrenTextLen = Array.from(el.children)
                .map(c => c.innerText.trim()).join('').replace(/\s+/g, '').length;
            const totalTextLen = totalText.replace(/\s+/g, '').length;
            // "has direct text" if total text is notably more than what children account for
            const hasDirectText = totalText && (
                el.children.length === 0 ||
                (allChildInline && totalTextLen > childrenTextLen + 1)
            );

            if (hasDirectText) {
                // First: preserve background/border styling (e.g. callout amber bg + left border)
                const bgImageLeaf = style.backgroundImage || 'none';
                const bgClipLeaf = style.webkitBackgroundClip || style.backgroundClip || '';
                // background-clip:text means gradient is used as text fill, not a visible background
                const hasGradientBgLeaf = bgImageLeaf !== 'none' && bgImageLeaf.includes('gradient') && bgClipLeaf !== 'text';
                if (hasBg || hasBorder || hasGradientBgLeaf) {
                    results.push({
                        type: 'shape', tag: tag, bounds: bounds,
                        styles: {
                            backgroundColor: bgColor,
                            backgroundImage: hasGradientBgLeaf ? bgImageLeaf : '',
                            border: borderStr, borderLeft: borderLeft, borderRight: borderRight,
                            borderTop: borderTop, borderBottom: borderBottom,
                            borderRadius: style.borderRadius
                        }
                    });
                }
                // Then: render the whole container as a text box (captures direct text + inline formatting)
                const {segments, gradientColors} = extractSegments(el);
                results.push({
                    type: 'text', tag: tag, text: totalText,
                    segments: segments, gradientColors: gradientColors,
                    textTransform: style.textTransform,
                    bounds: bounds,
                    styles: {
                        fontSize: style.fontSize, fontWeight: style.fontWeight,
                        fontFamily: style.fontFamily, letterSpacing: style.letterSpacing,
                        color: style.color, textAlign: style.textAlign,
                        lineHeight: style.lineHeight, listStyleType: style.listStyleType
                    }
                });
                return results;
            }

            // Standard container: maybe has background shape, then recurse
            const bgImage = style.backgroundImage || 'none';
            const hasGradientBg = bgImage !== 'none' && bgImage.includes('gradient');
            if (hasBg || hasBorder || hasGradientBg) {
                results.push({
                    type: 'shape',
                    tag: tag,
                    bounds: bounds,
                    styles: {
                        backgroundColor: bgColor,
                        backgroundImage: hasGradientBg ? bgImage : '',
                        border: borderStr,
                        borderLeft: borderLeft,
                        borderRight: borderRight,
                        borderTop: borderTop,
                        borderBottom: borderBottom,
                        borderRadius: style.borderRadius
                    }
                });
            }

            // Recurse into children
            for (const child of el.children) {
                results.push(...flatExtract(child));
            }
        }

        return results;
    }

    // Only process children of .slide-content (not the entire slide)
    const elements = [];
    for (const child of content.children) {
        elements.push(...flatExtract(child));
    }

    return {
        background: bgColor,
        elements: elements,
        slideSize: slideW && slideH ? {width: slideW/96, height: slideH/96} : null
    };
}
"""


def extract_slide_elements(page: Page, slide_index: int) -> Dict[str, Any]:
    result = page.evaluate(_EXTRACT_JS, slide_index)
    bg_rgb = parse_color(result['background']) if result['background'] else None
    return {
        'background': bg_rgb,
        'elements': result['elements'],
        'slideSize': result.get('slideSize')
    }


# CSS font → PPTX font mapping
# Key: substring to match in CSS fontFamily; Value: (latin_font, ea_font)
_FONT_MAP = {
    'Clash Display': ('Calibri Light', 'Microsoft YaHei'),
    'Satoshi':       ('Calibri',       'Microsoft YaHei'),
}
_DEFAULT_FONTS = ('Calibri', 'Microsoft YaHei')


def map_font(css_font_family: str):
    """Map CSS fontFamily string to (latin_font, ea_font) tuple."""
    if css_font_family:
        for css_name, fonts in _FONT_MAP.items():
            if css_name in css_font_family:
                return fonts
    return _DEFAULT_FONTS


def set_run_fonts(run, latin_font, ea_font):
    """Set latin, ea, cs fonts on a run element."""
    from lxml import etree
    NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    run.font.name = latin_font
    rPr = run._r.get_or_add_rPr()
    for tag, typeface in [('ea', ea_font), ('cs', ea_font)]:
        el = rPr.find(f'{{{NS}}}{tag}')
        if el is None:
            el = etree.SubElement(rPr, f'{{{NS}}}{tag}')
        el.set('typeface', typeface)


def set_letter_spacing(run, css_letter_spacing: str):
    """Set character spacing from CSS letterSpacing (e.g. '1.65px')."""
    if not css_letter_spacing or css_letter_spacing in ('normal', '0px'):
        return
    m = re.search(r'([\d.]+)px', css_letter_spacing)
    if m:
        px = float(m.group(1))
        # OOXML spc unit = 1/100 pt; 1px = 0.75pt → spc = px * 75
        spc = int(px * 75)
        if spc > 0:
            run._r.get_or_add_rPr().set('spc', str(spc))


def segments_to_lines(segments):
    """将 segments 列表拆分成行列表 (list of list)"""
    # Strip leading/trailing whitespace from text segments (HTML formatting whitespace)
    # But preserve \n (from BR tags) as line separators
    cleaned = []
    for s in segments:
        t = s['text']
        if t == '\n':
            cleaned.append(s)  # Keep newlines as-is
        elif t.strip():
            # Preserve original text (including leading/trailing spaces between inline elements
            # e.g. "<strong>Note:</strong> text" produces " text" which must keep its leading space)
            cleaned.append({'text': t, 'color': s['color'], 'bold': s.get('bold', False), 'fontSize': s.get('fontSize', '')})
    segments = cleaned
    lines = []
    current_line = []
    for seg in segments:
        text = seg['text']
        color = seg['color']
        bold = seg.get('bold', False)
        fontSize = seg.get('fontSize', '')
        if '\n' in text:
            parts = text.split('\n')
            for i, part in enumerate(parts):
                if part:
                    current_line.append({'text': part, 'color': color, 'bold': bold, 'fontSize': fontSize})
                if i < len(parts) - 1:
                    lines.append(current_line)
                    current_line = []
        else:
            current_line.append({'text': text, 'color': color, 'bold': bold, 'fontSize': fontSize})
    lines.append(current_line)
    # Strip leading/trailing empty lines, but preserve internal empty lines (from <BR><BR>)
    result = lines
    while result and not any(s['text'].strip() for s in result[0]):
        result = result[1:]
    while result and not any(s['text'].strip() for s in result[-1]):
        result = result[:-1]
    return result


def apply_run(run, text, color_str, font_size_pt, font_weight,
              text_transform='none', font_family='', letter_spacing=''):
    if text_transform == 'uppercase':
        text = text.upper()
    run.text = text
    # Preserve leading/trailing spaces in runs (OOXML strips them without xml:space="preserve")
    if text and (text[0] == ' ' or text[-1] == ' '):
        _nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        _t_elem = run._r.find('.//a:t', _nsmap)
        if _t_elem is not None:
            _t_elem.set('{http://www.w3.org/XML/1998/namespace}space', 'preserve')
    # P2: font mapping
    latin_font, ea_font = map_font(font_family)
    set_run_fonts(run, latin_font, ea_font)
    if font_size_pt:
        run.font.size = Pt(font_size_pt)
    try:
        if font_weight == 'bold':
            run.font.bold = True
        else:
            run.font.bold = int(font_weight) >= 600
    except Exception:
        pass
    rgb = parse_color(color_str)
    if rgb:
        run.font.color.rgb = RGBColor(*rgb)
    # P1: letter-spacing
    set_letter_spacing(run, letter_spacing)


def apply_para_format(p, s):
    lh = s.get('lineHeight', 'normal')
    if lh == 'normal':
        p.line_spacing = 1.2
    else:
        try:
            if 'px' in lh:
                lh_px = float(re.search(r'([\d.]+)', lh).group(1))
                # Use exact Pt line spacing so inline runs with larger fonts
                # don't inflate the line height beyond the CSS-measured value
                p.line_spacing = Pt(round(lh_px * 0.75, 1))
            else:
                p.line_spacing = float(lh)
        except Exception:
            p.line_spacing = 1.2
    align = s.get('textAlign', 'left')
    if align == 'center':
        p.alignment = PP_ALIGN.CENTER
    elif align == 'right':
        p.alignment = PP_ALIGN.RIGHT


def gradient_to_solid(bg_image, slide_bg=(13, 17, 23)):
    """Approximate a CSS gradient with a solid color by blending the first stop over slide bg."""
    if not bg_image or 'gradient' not in bg_image:
        return None
    rgba_matches = re.findall(r'rgba?\((\d+),\s*(\d+),\s*(\d+)(?:,\s*([\d.]+))?\)', bg_image)
    if not rgba_matches:
        return None
    r, g, b = int(rgba_matches[0][0]), int(rgba_matches[0][1]), int(rgba_matches[0][2])
    a = float(rgba_matches[0][3]) if rgba_matches[0][3] else 1.0
    if a <= 0:
        return None
    if a < 1.0:
        r = int(a * r + (1 - a) * slide_bg[0])
        g = int(a * g + (1 - a) * slide_bg[1])
        b = int(a * b + (1 - a) * slide_bg[2])
    return (r, g, b)


def export_shape_background(slide, elem, slide_bg=(255, 255, 255)):
    """仅创建背景形状（无文字），用于 type=shape 的容器"""
    b = elem['bounds']
    s = elem['styles']

    border_radius = s.get('borderRadius', '')
    radius_px = 0.0
    if border_radius and border_radius != '0px':
        m = re.search(r'([\d.]+)px', border_radius)
        if m:
            radius_px = float(m.group(1))

    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius_px > 0 else MSO_SHAPE.RECTANGLE,
        Inches(b['x']), Inches(b['y']),
        Inches(b['width']), Inches(b['height'])
    )
    if radius_px > 0:
        set_roundrect_adj(shape, radius_px, b['width'], b['height'])

    bg_rgb = parse_color(s.get('backgroundColor', ''), bg=slide_bg)
    if bg_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*bg_rgb)
    else:
        grad_fill = gradient_to_solid(s.get('backgroundImage', ''), slide_bg=slide_bg)
        if grad_fill:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*grad_fill)
        else:
            shape.fill.background()

    border_str = s.get('border', '')
    border_left = s.get('borderLeft', '')
    border_right = s.get('borderRight', '')
    border_top = s.get('borderTop', '')
    border_bottom = s.get('borderBottom', '')

    # Parse border sides
    def parse_border_side(bs):
        if not bs or 'none' in bs or bs.startswith('0px'):
            return None
        m = re.search(r'([\d.]+)px.*?rgb\((\d+),\s*(\d+),\s*(\d+)\)', bs)
        if m:
            return {'width': float(m.group(1)), 'rgb': (int(m.group(2)), int(m.group(3)), int(m.group(4)))}
        return None

    bl = parse_border_side(border_left)
    br = parse_border_side(border_right)
    bt = parse_border_side(border_top)
    bb = parse_border_side(border_bottom)

    sides_with_border = sum(1 for x in [bl, br, bt, bb] if x is not None)
    borders = [x for x in [bl, br, bt, bb] if x is not None]
    all_uniform = (len(borders) >= 3 and
                   all(bd['rgb'] == borders[0]['rgb'] and bd['width'] == borders[0]['width']
                       for bd in borders))

    if all_uniform:
        # All borders same color/width → use shape.line
        shape.line.color.rgb = RGBColor(*borders[0]['rgb'])
        shape.line.width = Pt(max(0.5, borders[0]['width'] * 0.75))
    elif sides_with_border >= 1:
        # Mixed or partial borders → draw each side as separate rectangle
        suppress_line(shape)
        if bl:
            border_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(b['x']), Inches(b['y']),
                Inches(bl['width'] / 96), Inches(b['height'])
            )
            border_shape.fill.solid()
            border_shape.fill.fore_color.rgb = RGBColor(*bl['rgb'])
            suppress_line(border_shape)
        if br:
            border_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(b['x'] + b['width'] - br['width']/96), Inches(b['y']),
                Inches(br['width'] / 96), Inches(b['height'])
            )
            border_shape.fill.solid()
            border_shape.fill.fore_color.rgb = RGBColor(*br['rgb'])
            suppress_line(border_shape)
        if bt:
            border_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(b['x']), Inches(b['y']),
                Inches(b['width']), Inches(bt['width'] / 96)
            )
            border_shape.fill.solid()
            border_shape.fill.fore_color.rgb = RGBColor(*bt['rgb'])
            suppress_line(border_shape)
        if bb:
            border_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(b['x']), Inches(b['y'] + b['height'] - bb['width']/96),
                Inches(b['width']), Inches(bb['width'] / 96)
            )
            border_shape.fill.solid()
            border_shape.fill.fore_color.rgb = RGBColor(*bb['rgb'])
            suppress_line(border_shape)
    else:
        # No border
        suppress_line(shape)

    # Remove default text frame content
    tf = shape.text_frame
    if tf.paragraphs:
        for para in tf.paragraphs:
            for run in para.runs:
                run.text = ''

    return shape


def interpolate_color(c1, c2, t):
    """在两色间插值，t=0返回c1，t=1返回c2"""
    return (
        int(c1[0] + (c2[0] - c1[0]) * t),
        int(c1[1] + (c2[1] - c1[1]) * t),
        int(c1[2] + (c2[2] - c1[2]) * t)
    )


def export_text_element(slide, elem: Dict[str, Any], bg_color=None):
    b = elem['bounds']
    s = elem['styles']
    segments = elem.get('segments', [])
    text_transform = elem.get('textTransform', 'none')
    font_size_pt = px_to_pt(s.get('fontSize', '16px'))
    font_weight = s.get('fontWeight', '400')
    font_family = s.get('fontFamily', '')
    letter_spacing = s.get('letterSpacing', '')

    txBox = slide.shapes.add_textbox(
        Inches(b['x']), Inches(b['y']),
        Inches(b['width']), Inches(b['height'])
    )
    tf = txBox.text_frame
    # Disable word wrap for very large decorative text (e.g. chapter numbers at 100+pt)
    # to prevent character-level line breaking in narrow boxes
    tf.word_wrap = font_size_pt < 60
    from pptx.enum.text import MSO_AUTO_SIZE
    tf.auto_size = MSO_AUTO_SIZE.NONE

    if not segments:
        raw = (elem.get('text', '') or '').strip()
        segments = [{'text': raw, 'color': s.get('color', '')}]

    lines = segments_to_lines(segments)
    if not lines:
        lines = [[{'text': '', 'color': s.get('color', '')}]]

    # H1 渐变近似：多行时按比例分配渐变色
    gradient_colors = elem.get('gradientColors') if elem.get('tag') == 'h1' else None
    gc_start = parse_color(gradient_colors[0]) if gradient_colors else None
    gc_end = parse_color(gradient_colors[1]) if gradient_colors else None
    total_lines = len(lines)

    # li 元素：添加蓝色 ▶ 前缀
    is_li = elem.get('tag') == 'li'
    li_bullet_color = 'rgb(56, 139, 253)'  # --accent-blue

    for line_idx, line_segs in enumerate(lines):
        p = tf.add_paragraph() if line_idx > 0 else tf.paragraphs[0]
        apply_para_format(p, s)

        # li 前缀（只在第一行加）
        if is_li and line_idx == 0:
            bullet_run = p.add_run()
            apply_run(bullet_run, '▶ ', li_bullet_color, font_size_pt * 0.7, '400')

        # 计算当前行的渐变色
        if gc_start and gc_end and total_lines > 1:
            t = line_idx / (total_lines - 1)
            grad_rgb = interpolate_color(gc_start, gc_end, t)
            override_color = 'rgb({},{},{})'.format(*grad_rgb)
        elif gc_start and gc_end:
            override_color = gradient_colors[1]
        else:
            override_color = None

        for seg in line_segs:
            if not seg['text']:
                continue
            run = p.add_run()
            # 如果有渐变色覆盖，使用覆盖色；否则用 seg 原色
            color = override_color or seg['color']
            # 使用 segment 级别的 bold（如果有），否则用元素级别
            seg_weight = 'bold' if seg.get('bold') else font_weight
            seg_fs_raw = seg.get('fontSize', '')
            seg_font_size_pt = px_to_pt(seg_fs_raw) if seg_fs_raw and 'px' in str(seg_fs_raw) else font_size_pt
            apply_run(run, seg['text'], color, seg_font_size_pt, seg_weight, text_transform,
                      font_family=font_family, letter_spacing=letter_spacing)


def export_shape_with_text(slide, elem: Dict[str, Any], bg_color=None):
    b = elem['bounds']
    s = elem['styles']

    border_radius = s.get('borderRadius', '')
    radius_px = 0.0
    if border_radius and border_radius != '0px':
        m = re.search(r'([\d.]+)px', border_radius)
        if m:
            radius_px = float(m.group(1))

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius_px > 0 else MSO_SHAPE.RECTANGLE,
        Inches(b['x']), Inches(b['y']),
        Inches(b['width']), Inches(b['height'])
    )
    if radius_px > 0:
        set_roundrect_adj(shape, radius_px, b['width'], b['height'])

    bg_rgb = parse_color(s.get('backgroundColor', ''))
    if bg_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*bg_rgb)
    else:
        grad_fill = gradient_to_solid(s.get('backgroundImage', ''))
        if grad_fill:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*grad_fill)
        else:
            shape.fill.background()

    border_str = s.get('border', '')
    if border_str and 'none' not in border_str:
        m = re.search(r'rgb\((\d+),\s*(\d+),\s*(\d+)\)', border_str)
        if m:
            shape.line.color.rgb = RGBColor(int(m.group(1)), int(m.group(2)), int(m.group(3)))
    else:
        shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(10)
    tf.margin_bottom = Pt(10)

    segments = elem.get('segments', [])
    font_size_pt = px_to_pt(s.get('fontSize', '16px'))
    font_weight = s.get('fontWeight', '400')
    font_family = s.get('fontFamily', '')
    letter_spacing = s.get('letterSpacing', '')
    text_transform = elem.get('textTransform', 'none')

    if not segments:
        raw = (elem.get('text', '') or '').strip()
        segments = [{'text': raw, 'color': s.get('color', '')}]

    lines = segments_to_lines(segments)
    for line_idx, line_segs in enumerate(lines):
        p = tf.add_paragraph() if line_idx > 0 else tf.paragraphs[0]
        apply_para_format(p, s)
        for seg in line_segs:
            if not seg['text']:
                continue
            run = p.add_run()
            seg_weight = 'bold' if seg.get('bold') else font_weight
            apply_run(run, seg['text'], seg['color'], font_size_pt, seg_weight, text_transform,
                      font_family=font_family, letter_spacing=letter_spacing)


def export_table_element(slide, elem: Dict[str, Any]):
    """Render HTML table as individual cell rectangles + text frames."""
    rows = elem.get('rows', [])
    if not rows:
        return

    for row_data in rows:
        for cell in row_data['cells']:
            cb = cell['bounds']
            cs = cell['styles']

            # Skip zero-size cells
            if cb['width'] < 0.01 or cb['height'] < 0.01:
                continue

            # Cell background rectangle
            cell_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(cb['x']), Inches(cb['y']),
                Inches(cb['width']), Inches(cb['height'])
            )

            bg_rgb = parse_color(cs.get('backgroundColor', ''))
            if bg_rgb:
                cell_shape.fill.solid()
                cell_shape.fill.fore_color.rgb = RGBColor(*bg_rgb)
            else:
                cell_shape.fill.background()

            # No rectangle border (avoids vertical dividers)
            suppress_line(cell_shape)

            # border-bottom only: draw as a separate thin rectangle
            border_bottom = cs.get('borderBottom', '')
            if border_bottom and 'none' not in border_bottom and not border_bottom.startswith('0px'):
                m = re.search(r'rgb\((\d+),\s*(\d+),\s*(\d+)\)', border_bottom)
                if m:
                    divider = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(cb['x']), Inches(cb['y'] + cb['height'] - 0.005),
                        Inches(cb['width']), Inches(0.005)
                    )
                    divider.fill.solid()
                    divider.fill.fore_color.rgb = RGBColor(int(m.group(1)), int(m.group(2)), int(m.group(3)))
                    suppress_line(divider)

            # Text content
            segments = cell.get('segments', [])
            text = cell.get('text', '').strip()
            if not segments and text:
                segments = [{'text': text, 'color': cs.get('color', '')}]
            if not segments:
                continue

            font_size_pt = px_to_pt(cs.get('fontSize', '14px'))
            font_weight = cs.get('fontWeight', '400')
            font_family = cs.get('fontFamily', '')
            letter_spacing = cs.get('letterSpacing', '')
            if cell['isHeader']:
                font_weight = 'bold'

            tf = cell_shape.text_frame
            tf.word_wrap = True
            tf.margin_left = Pt(6)
            tf.margin_right = Pt(6)
            tf.margin_top = Pt(4)
            tf.margin_bottom = Pt(4)

            lines = segments_to_lines(segments)
            for line_idx, line_segs in enumerate(lines):
                p = tf.add_paragraph() if line_idx > 0 else tf.paragraphs[0]
                align = cs.get('textAlign', 'left')
                if align == 'center':
                    p.alignment = PP_ALIGN.CENTER
                elif align == 'right':
                    p.alignment = PP_ALIGN.RIGHT
                for seg in line_segs:
                    if not seg['text']:
                        continue
                    run = p.add_run()
                    seg_weight = 'bold' if seg.get('bold') else font_weight
                    apply_run(run, seg['text'], seg['color'], font_size_pt, seg_weight,
                              font_family=font_family, letter_spacing=letter_spacing)


def export_native(html_path, output_path=None, width=1440, height=900):
    html_path = Path(html_path).resolve()
    if not html_path.exists():
        print(f"Error: {html_path}")
        sys.exit(1)

    output_path = Path(output_path) if output_path else html_path.with_suffix('.pptx')
    print(f"导出（native v4）: {html_path.name}")

    with sync_playwright() as p:
        browser = p.chromium.launch(channel='chrome', headless=True)
        page = browser.new_page(viewport={"width": width, "height": height})
        page.goto(f"file://{html_path}", wait_until="networkidle")
        page.wait_for_timeout(500)

        slide_count = page.evaluate("document.querySelectorAll('.slide').length")
        if slide_count == 0:
            print("未找到 .slide 元素")
            browser.close()
            return

        print(f"找到 {slide_count} 张幻灯片")
        inject_visible(page)
        page.wait_for_timeout(200)

        prs = Presentation()
        first = extract_slide_elements(page, 0)
        if first.get('slideSize'):
            prs.slide_width = Inches(first['slideSize']['width'])
            prs.slide_height = Inches(first['slideSize']['height'])
        else:
            prs.slide_width = Inches(15.0)
            prs.slide_height = Inches(9.375)

        blank_layout = prs.slide_layouts[6]
        max_w = prs.slide_width / 914400

        for i in range(slide_count):
            print(f"  [{i+1}/{slide_count}] 处理中...")
            data = extract_slide_elements(page, i)
            slide = prs.slides.add_slide(blank_layout)

            if data['background']:
                r, g, b = data['background']
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = RGBColor(r, g, b)

            for elem in data['elements']:
                try:
                    # Clamp width
                    if elem['bounds']['x'] < max_w and elem['bounds']['width'] > max_w - elem['bounds']['x']:
                        elem['bounds']['width'] = max_w - elem['bounds']['x']

                    elem_type = elem.get('type', 'text')

                    if elem_type == 'shape':
                        export_shape_background(slide, elem, slide_bg=data['background'] or (255, 255, 255))
                    elif elem_type == 'table':
                        export_table_element(slide, elem)
                    else:
                        export_text_element(slide, elem, data['background'])
                except Exception as e:
                    print(f"    警告: {e}")

        browser.close()

    prs.save(str(output_path))
    print(f"✓ 已保存: {output_path}  ({slide_count} 张幻灯片)")
    return output_path


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("html")
    parser.add_argument("output", nargs="?")
    parser.add_argument("--width", type=int, default=1440)
    parser.add_argument("--height", type=int, default=900)
    args = parser.parse_args()
    export_native(args.html, args.output, args.width, args.height)


if __name__ == "__main__":
    main()
